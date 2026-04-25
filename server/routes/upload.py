"""
routes/upload.py — document upload + text extraction.

POST /v1/upload
───────────────
Accepts a single file (pdf, docx, pptx, xlsx, xls) up to 10 MB.
Returns the extracted plain text so the frontend can inject it into a
chat message without needing another round-trip when the user hits send.
"""

import io
import logging
import os

from fastapi import APIRouter, Depends, File, HTTPException, UploadFile, status
from fastapi.security import HTTPAuthorizationCredentials, HTTPBearer
from pydantic import BaseModel

from security.security import decode_access_token

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/v1", tags=["upload"])
_bearer = HTTPBearer()

_MAX_BYTES = 10 * 1024 * 1024  # 10 MB
_ALLOWED_EXT = {".pdf", ".docx", ".pptx", ".xlsx", ".xls"}


def _get_current_user_id(
    credentials: HTTPAuthorizationCredentials = Depends(_bearer),
) -> int:
    payload = decode_access_token(credentials.credentials)
    try:
        return int(payload["id"])
    except (KeyError, ValueError):
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Invalid token")


# ── extractors ─────────────────────────────────────────────────────────────────


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    texts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if slide_texts:
            texts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
    return "\n\n".join(texts)


def _extract_xlsx(data: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    rows: list[str] = []
    for sheet in wb.worksheets:
        rows.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append("\t".join(cells))
    wb.close()
    return "\n".join(rows)


# ── schema ─────────────────────────────────────────────────────────────────────


class UploadResponse(BaseModel):
    filename: str
    file_type: str
    text: str
    char_count: int


# ── endpoint ───────────────────────────────────────────────────────────────────


@router.post("/upload", response_model=UploadResponse)
async def upload_document(
    file: UploadFile = File(...),
    user_id: int = Depends(_get_current_user_id),
) -> UploadResponse:
    """Extract plain text from an uploaded document (PDF, DOCX, PPTX, XLSX/XLS)."""
    ext = os.path.splitext(file.filename or "")[1].lower()
    if ext not in _ALLOWED_EXT:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Unsupported file type '{ext}'. Allowed: {', '.join(sorted(_ALLOWED_EXT))}",
        )

    data = await file.read()
    if len(data) > _MAX_BYTES:
        raise HTTPException(
            status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
            detail="File exceeds the 10 MB limit.",
        )

    try:
        if ext == ".pdf":
            text = _extract_pdf(data)
        elif ext == ".docx":
            text = _extract_docx(data)
        elif ext == ".pptx":
            text = _extract_pptx(data)
        else:
            text = _extract_xlsx(data)
    except Exception as exc:
        logger.exception("Text extraction failed for %s", file.filename)
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"Could not extract text: {exc}",
        )

    text = text.strip()
    logger.info("[UPLOAD] user_id=%d  file=%s  chars=%d", user_id, file.filename, len(text))

    return UploadResponse(
        filename=file.filename or "document",
        file_type=ext.lstrip("."),
        text=text,
        char_count=len(text),
    )
